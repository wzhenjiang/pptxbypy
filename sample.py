#! /usr/local/bin/python3

# example file for powerpoint slides generation
import random
import os
from pptx import Presentation
from pptx.util import Inches

# if font operation required, pls include below import statement
from openpyxl.styles import Font

class WorkBook():
	''' WorkBook is to wrap the workbook. each instance is a workbook
	'''

	def __init__(self,filename=None):
		''' if filename not provided, a new workbook would be created
		'''
		self.status = False

		if not filename:
			self.filename = "temp.xls.file.%04d.xlsx" % random.randint(1,1000) 
			self.workbook = openpyxl.Workbook()
			self.workbook.create_sheet(title="Sheet1", index=0)
			self.workbook.save(self.filename)
		else:
			self.filename = filename
		self.load_excel()

	def load_excel(self):
		''' load excel sheets into the instance
		'''
		try:
			self.workbook = openpyxl.load_workbook(self.filename)
			self.status = True
			print('[%s] loaded. %d sheets inside: %s' % (self.filename , len(self.workbook.sheetnames),self.workbook.sheetnames))
		except FileNotFoundError:
			print('Load error: %s not found' % self.filename)

	def get_sheet(self, index=0):
		if not self.status:
			return
		self.sheet = self.workbook[self.workbook.sheetnames[index]]
		print('Sheet loaded, data range: [%d:%d]' % (self.sheet.max_row,self.sheet.max_column))
		print('Sammple value in [A1 / B3] : [%s / %s]' % ( self.sheet['A1'].value, self.sheet['B3'].value))

	def manipulate_and_save(self):
		''' sample code of data manipulation
		'''
		print("SET TITLE / DATA; CREATE NEW SHEET; SET DIMENSION AND FONT")
		self.sheet.title = 'SAMPLE TITLE'
		self.sheet['A1'] = 1999
		self.sheet['B3'] = 'aloha'
		for i in range(8,21):
			for j in range(1,11):
				self.sheet.cell(row = i, column = j).value = i*10+j
		self.workbook.create_sheet(title='new sheet' ,index = 1)
		self.sheet.column_dimensions['B'].width = 40
		self.sheet.row_dimensions[1].height = 30
		self.sheet['B3'].font = Font(sz=18,bold=True, italic=True)
		self.workbook.save(self.filename + '.out.xlsx')

def fill_deliverable_page(shapes):
	shapes.title.text= "Deliverable Tracking"
	cats = ("Deliverables",
		"Issue / Concern",
		"Mitigation Plan",
		"Action / Key Ask",
		"Owner",
		"Due Date",
		"S",
		"Status")
	widths = (1.5,2.5,1.5,2,1.5,1.5,0.3,1)

	row = 10 
	col = len(cats)
	
	left = Inches(1)
	top = Inches(1.3)
	width = Inches(7.0)
	height = Inches(0.8)
	
	table = shapes.add_table(row,col,left,top,width,height).table

	i = 0
	for cat in cats:
		table.cell(0,i).text = cat
		table.columns[i].width = Inches(widths[i])
		i += 1

# unit test code below
if __name__ == '__main__':

	#curdir = os.path.dirname(os.path.realpath('.'))
	#print("Current Dir: %s" % curdir)

	# first, new an presentation object
	prs = Presentation('template.pptx')

	# second, get layout and use it to generate one page
	cover_page_layout = prs.slide_layouts[0]
	cover_page= prs.slides.add_slide(cover_page_layout)

	deliverable_page_layout = prs.slide_layouts[2]
	deliverable_page= prs.slides.add_slide(deliverable_page_layout)

	print("%d placeholders inside" % len(cover_page.placeholders))
	
	for shape in cover_page.placeholders:
		print('%d %s' % (shape.placeholder_format.idx, shape.name))
	
	print("%d shapes inside" % len(deliverable_page.shapes))
	
	for shape in deliverable_page.placeholders:
		print('%d %s' % (shape.placeholder_format.idx, shape.name))

	table = fill_deliverable_page(deliverable_page.shapes)
	# third, set title and subtitle
	cover_page.placeholders[13].text = "5GTF STREAM-1 SUBSTREAM-2 CORE UPSTREAM WEEKLY" 
	cover_page.placeholders[14].text = "by xBG solution team led by MN PS"

	prs.save('test.pptx')

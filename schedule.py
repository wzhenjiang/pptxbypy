#! /usr/local/bin/python3

# example file for powerpoint slides generation
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches                # import Inches as it is most common
from pptx.util import Cm                    # import Cm as it is what used in powerpoint windows
from pptx.dml.color import RGBColor         # import color for color definition
from pptx.enum.shapes import MSO_SHAPE      # import auto shape for shape operation

# more modules as reference below
from pptx.enum.dml import MSO_THEME_COLOR
# chart data related
# from pptx.chart.data import ChartData
# from pptx.enum.chart import XL_TICK_MARK
# from pptx.enum.chart import XL_DATA_LABEL_POSITION
# from pptx.enum.chart import XL_LEGEND_POSITON
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.enum.chart import XL_MARKER_STYLE

# define variables

# define page content area size 1 inch is 2.54 cm or 79 Pts
CONTENT_AREA_WIDTH = Cm(30)     # 12 inches
CONTENT_AREA_HEIGHT = Cm(15)     # 6 inches 

# define layout index from slidemaster count
COVER_PAGE_LAYOUT = 0
DELIVERY_PAGE_LAYOUT = 1
CONTENT_PAGE_LAYOUT = 2
SCHEDULE_LAYOUT = 5

# define chevron shape size
CHEVRON_WIDTH = Cm(2)
CHEVRON_HEIGHT = Cm(0.6)

# define CHERON shape color for schedule geeration
SCHEDULE_FILL_COLOR = RGBColor(0,0,128)
SCHEDULE_LINE_COLOR = RGBColor(200,200,200)
SCHEDULE_LINE_WIDTH = Pt(1)

def create_schedule_frame(slide,pos,start,end,title):
    ''' this fucntion add schedule shapes
    '''
    left = Inches(pos[0])
    top = Inches(pos[1])
    shapes = slide.shapes
    for x in range(start,end+1):
        print("Generating",x,"blocks...")
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON,left, top, CHEVRON_WIDTH, CHEVRON_HEIGHT)                       
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = SCHEDULE_FILL_COLOR
        line = shape.line
        line.color.rgb = SCHEDULE_LINE_COLOR
        line.width = SCHEDULE_LINE_WIDTH
        shape.text = title + str(x)
        left += CHEVRON_WIDTH


if __name__ == '__main__':
    '''                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.clear()
                shape.text = "P02"
    '''
    prs = Presentation('template.pptx')
    slide_layout = prs.slide_layouts[SCHEDULE_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    create_schedule_frame(slide,[1,1],2,10,'P')
    prs.save('schedule.pptx')

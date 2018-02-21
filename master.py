#! /usr/local/bin/python3

# example file for powerpoint slides generation
from pptx import Presentation
from pptx.util import Inches

# more modules as reference below
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# chart data related
# from pptx.chart.data import ChartData
# from pptx.enum.chart import XL_TICK_MARK
# from pptx.enum.chart import XL_DATA_LABEL_POSITION
# from pptx.enum.chart import XL_LEGEND_POSITON
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.enum.chart import XL_MARKER_STYLE

# define variables
COVER_PAGE_LAYOUT = 0
DELIVERY_PAGE_LAYOUT = 1
CONTENT_PAGE_LAYOUT = 2
COVER_TITLE_IDX = 13
COVER_SUBTITLE_IDX = 14

class AutoSlide():
	''' class to wrap slides related function
	'''
	def __init__(self, filename = None):
		if filename:
			self.templatefn = filename
		else:
			self.templatefn = "template.pptx"
		self.presentation = Presentation(self.templatefn)
		print("Initiate presentation w/ %s" % self.templatefn)
		self.layouts = self.presentation.slide_layouts
		self.slides = self.presentation.slides

	def add_cover_page(self, title, subtitle):
		cover_page_layout = self.layouts[COVER_PAGE_LAYOUT]
		self.coverpage = self.slides.add_slide(cover_page_layout)
		for shape in self.coverpage.placeholders:
			print('<DEBUG placeholder> %d:%s' % (shape.placeholder_format.idx, shape.name))
		# third, set title and subtitle
		self.coverpage.placeholders[COVER_TITLE_IDX].text = title 
		self.coverpage.placeholders[COVER_SUBTITLE_IDX].text = subtitle
		self.coverpage.name = "Cover"

	def add_delivery_page(self, page_title, column_titles, column_widths):
		# add the page
		self.delivery_page = self.slides.add_slide(self.layouts[DELIVERY_PAGE_LAYOUT])
		self.delivery_page.name = "Delivery"

		# set page title
		shapes = self.delivery_page.shapes
		shapes.title.text= page_title 
		
		# create table
		num_row = 10 
		num_col = len(column_titles)
	
		left = Inches(1)
		top = Inches(1.3)
		width = Inches(7.0)
		height = Inches(0.8)
	
		table = shapes.add_table(num_row,num_col,left,top,width,height).table

		# set value for table titles
		i = 0
		for title in column_titles:
			table.cell(0,i).text = title
			table.cell(0,i).fill.solid()
			table.cell(0,i).fill.fore_color.rgb = RGBColor(0x0,0x0,0x5f)
			table.cell(0,i).text_frame.paragraphs[0].font.size = Pt(12)
			table.columns[i].width = Inches(column_widths[i])
			i += 1
		table.cell(1,len(column_titles)-2).fill.solid()
		table.cell(1,len(column_titles)-2).fill.fore_color.rgb = RGBColor(0xff,0x0,0x0)
		
# unit test code below
if __name__ == '__main__':

	# new an presentation object
	presentation = AutoSlide('template.pptx')

	# prapare data for cover page
	title = "5GTF STREAM-1 SUBSTREAM-2 CORE UPSTREAM WEEKLY" 
	subtitle = "by xBG solution team led by MN PS"
	
	# Generate cover page
	presentation.add_cover_page(title,subtitle)

	# prepare dat for deliverable page
	page_title = "Deliverable Tracking"
	column_titles = ("Deliverables",
			"Issue / Concern",
			"Mitigation Plan",
			"Action / Key Ask",
			"Owner",
			"Due Date",
			"S",
			"Status")
	column_widths = (1.5,2.5,1.5,2,1.5,1.5,0.3,1)

	# generate deliverable taracking page
	presentation.add_delivery_page(page_title,column_titles, column_widths)

	# save the generated file
	presentation.presentation.save('test.pptx')
